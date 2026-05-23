import customtkinter as ctk
import pandas as pd
from tkinter import messagebox, ttk, filedialog
import re
import os
import glob # Dosya arama motoru için eklendi

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    pass

ctk.set_appearance_mode("Dark")
ctk.set_default_color_theme("blue")

PREGATE_EKIBI = ["tolga tosun", "muhammt acı", "serhat seymen", "muhammet malik eşber kılıç", "emre sabe", "kaan turgut", "oğuzhan öz"]

KATEGORILER = [
    "Sorumluluk Evrakı", "Yüksekte Çalışabilir Sağlık Raporu", "Tank Basınç Raporu - T9",
    "Muayene", "Zorunlu Trafik Sigortası", "K1 Taşıt Kartı", "Tehlikeli Atık Sigortası",
    "Tank Kodu", "KKD", "interchange", "Şoför Kart Almaya Gelmedi",
    "Sevkiyat - Firma Sürücü İptal", "Havuz Tahliye Vanası", "İSOPA",
    "FAR-STOP-SİS FARI-LASTİK KIRIK/PATLAK", "Turuncu Plaka ve İkaz Tabelası",
    "Yakıt Deposu Sızıntı", "Diğer"
]

class PregateDashboard(ctk.CTk):
    def __init__(self):
        super().__init__()
        self.title("Pregate Operasyon Dashboard - S. SEYMEN")
        self.geometry("1300x850")
        self.df_dinamik = pd.DataFrame()
        self.df_pregate = pd.DataFrame()
        self.df_pol_temp = pd.DataFrame()
        
        self.grid_rowconfigure(0, weight=1)
        self.grid_columnconfigure(1, weight=1)
        
        self.arayuzu_olustur()

    def arayuzu_olustur(self):
        style = ttk.Style()
        style.theme_use("default")
        style.configure("Treeview", background="#2b2b2b", foreground="white", rowheight=30, fieldbackground="#2b2b2b", borderwidth=0, font=("Arial", 11))
        style.map('Treeview', background=[('selected', '#347083')])
        style.configure("Treeview.Heading", background="#1f242d", foreground="#F1C40F", font=("Arial", 11, "bold"), relief="flat")
        style.map("Treeview.Heading", background=[('active', '#3484F0')])
        
        # --- SOL MENÜ ---
        self.sidebar = ctk.CTkFrame(self, width=280, corner_radius=0)
        self.sidebar.grid(row=0, column=0, sticky="nsew")

        ctk.CTkLabel(self.sidebar, text="Menü", font=("Arial", 26, "bold")).pack(pady=(15, 10))
        
        self.btn_dinamik_yukle = ctk.CTkButton(self.sidebar, text="Dinamik Rapor Yükle", command=self.dinamik_yukle)
        self.btn_dinamik_yukle.pack(pady=5, padx=20)
        self.btn_pregate_yukle = ctk.CTkButton(self.sidebar, text="Pregate Sunum Yükle", command=self.pregate_yukle)
        self.btn_pregate_yukle.pack(pady=5, padx=20)

        ctk.CTkLabel(self.sidebar, text="Raporlama Çıktıları", font=("Arial", 14, "bold"), text_color="#F1C40F").pack(pady=(15, 5))
        self.btn_excel = ctk.CTkButton(self.sidebar, text="📊 Excel Olarak Kaydet", fg_color="#D4AC0D", hover_color="#9A7D0A", text_color="black", command=self.export_excel)
        self.btn_excel.pack(pady=5, padx=20, fill="x")
        self.btn_pptx = ctk.CTkButton(self.sidebar, text="🎬 Sunum (PPTX) Oluştur", fg_color="#D35400", hover_color="#A04000", command=self.export_pptx)
        self.btn_pptx.pack(pady=5, padx=20, fill="x")

        self.btn_anasayfa = ctk.CTkButton(self.sidebar, text="🏠 ANA SAYFA", fg_color="#e67e22", hover_color="#d35400", command=self.ana_sayfayi_goster)
        self.btn_anasayfa.pack(pady=(15, 5), padx=20, fill="x", ipady=5)
        
        ctk.CTkLabel(self.sidebar, text="Dinamik Rapor (Fabrika)", font=("Arial", 14, "bold"), text_color="#E74C3C").pack(pady=(10, 5))
        for dep in ["Kimya", "Dow", "Diğer"]:
            ctk.CTkButton(self.sidebar, text=dep, fg_color="#922B21", hover_color="#641E16", command=lambda d=dep: self.dinamik_rapor_analiz(d)).pack(pady=3, padx=20, fill="x")
            
        ctk.CTkLabel(self.sidebar, text="Pregate (Terminal)", font=("Arial", 14, "bold"), text_color="#3498DB").pack(pady=(10, 5))
        for dep in ["Poliport Terminal", "Antrepo"]:
            ctk.CTkButton(self.sidebar, text=dep, fg_color="#1A5276", hover_color="#154360", command=lambda d=dep: self.pregate_analiz(d)).pack(pady=3, padx=20, fill="x")
            
        self.btn_nakliyeci = ctk.CTkButton(self.sidebar, text="🚛 Poliport Nakliyeciler", fg_color="#2874A6", hover_color="#1B4F72", font=("Arial", 14, "bold"), command=self.poliport_nakliyeci_arayuz)
        self.btn_nakliyeci.pack(pady=(15, 5), padx=20, fill="x", ipady=5)

        # =========================================================
        # İMZA (MARKA) ALANI - SOL ALT SABİT
        # =========================================================
        self.lbl_imza = ctk.CTkLabel(
            self.sidebar, 
            text="S. SEYMEN", 
            font=("Lucida Handwriting", 24, "italic"),
            text_color="#D4AF37" 
        )
        self.lbl_imza.pack(side="bottom", pady=25)
            
        # --- ANA EKRAN ---
        self.main_frame = ctk.CTkFrame(self, fg_color="transparent")
        self.main_frame.grid(row=0, column=1, sticky="nsew", padx=20, pady=20)
        
        self.lbl_secili_dep = ctk.CTkLabel(self.main_frame, text="Sisteme Hoş Geldiniz. Lütfen Verileri Yükleyin.", font=("Arial", 26, "bold"))
        self.lbl_secili_dep.pack(pady=(0, 15))
        
        self.kpi_frame = ctk.CTkFrame(self.main_frame, fg_color="transparent")
        self.kpi_frame.pack(fill="x", pady=10)
        
        self.lbl_toplam = self.kpi_kart_olustur(self.kpi_frame, "Toplam Araç", "0", "#3498db", 0)
        self.lbl_onay = self.kpi_kart_olustur(self.kpi_frame, "Onaylanan", "0", "#2ecc71", 1)
        self.lbl_red = self.kpi_kart_olustur(self.kpi_frame, "Reddedilen", "0", "#e74c3c", 2)
        
        self.nakliye_menu = ctk.CTkFrame(self.main_frame, fg_color="transparent")
        self.btn_n_tum = ctk.CTkButton(self.nakliye_menu, text="Tüm Nakliyeci Kayıtları", fg_color="#34495E", command=self.n_listele_tum)
        self.btn_n_tum.pack(side="left", padx=10, expand=True, fill="x")
        self.btn_n_onay = ctk.CTkButton(self.nakliye_menu, text="Onaylanan Nakliyeciler", fg_color="#27AE60", command=self.n_listele_onay)
        self.btn_n_onay.pack(side="left", padx=10, expand=True, fill="x")
        self.btn_n_red = ctk.CTkButton(self.nakliye_menu, text="Reddedilenler & Matrix Tablosu", fg_color="#C0392B", command=self.n_listele_red)
        self.btn_n_red.pack(side="left", padx=10, expand=True, fill="x")
        
        self.dashboard_frame = ctk.CTkFrame(self.main_frame, fg_color="transparent")
        self.txt_rapor = ctk.CTkTextbox(self.main_frame, font=("Consolas", 15))
        
        self.tree_frame = ctk.CTkFrame(self.main_frame, fg_color="transparent")
        self.tree_scroll_y = ttk.Scrollbar(self.tree_frame, orient="vertical")
        self.tree_scroll_y.pack(side="right", fill="y")
        self.tree_scroll_x = ttk.Scrollbar(self.tree_frame, orient="horizontal")
        self.tree_scroll_x.pack(side="bottom", fill="x")
        self.tree = ttk.Treeview(self.tree_frame, yscrollcommand=self.tree_scroll_y.set, xscrollcommand=self.tree_scroll_x.set)
        self.tree.pack(side="left", fill="both", expand=True)
        self.tree_scroll_y.configure(command=self.tree.yview)
        self.tree_scroll_x.configure(command=self.tree.xview)

        # TEXTBOX STİLLERİ
        self.txt_rapor._textbox.tag_configure("baslik", background="#1F242D", foreground="#F1C40F", font=("Arial", 18, "bold"), justify="left", spacing1=15, spacing3=15)
        self.txt_rapor._textbox.tag_configure("altbaslik", font=("Arial", 15, "bold"), foreground="#3498DB", spacing1=10, spacing3=5)
        self.txt_rapor._textbox.tag_configure("satir_tek", background="#212F3C", foreground="#ECF0F1", font=("Consolas", 15), spacing1=6, spacing3=6, tabs=(480, 630, 780))
        self.txt_rapor._textbox.tag_configure("satir_cift", background="#2C3E50", foreground="#ECF0F1", font=("Consolas", 15), spacing1=6, spacing3=6, tabs=(480, 630, 780))

    def kpi_kart_olustur(self, parent, baslik, deger, renk, sutun):
        kart = ctk.CTkFrame(parent, corner_radius=10)
        kart.grid(row=0, column=sutun, padx=10, sticky="nsew")
        parent.grid_columnconfigure(sutun, weight=1)
        ctk.CTkLabel(kart, text=baslik, font=("Arial", 18)).pack(pady=(15, 0))
        lbl_deger = ctk.CTkLabel(kart, text=deger, font=("Arial", 36, "bold"), text_color=renk)
        lbl_deger.pack(pady=(5, 15))
        return lbl_deger

    def mini_karti_olustur(self, parent, baslik, to, on, re_val, r_satir, r_sutun):
        kart = ctk.CTkFrame(parent, corner_radius=8, border_width=1, border_color="#34495e")
        kart.grid(row=r_satir, column=r_sutun, padx=10, pady=10, sticky="nsew")
        parent.grid_columnconfigure(r_sutun, weight=1)
        ctk.CTkLabel(kart, text=baslik, font=("Arial", 16, "bold")).pack(pady=(10, 5))
        info_frame = ctk.CTkFrame(kart, fg_color="transparent")
        info_frame.pack(fill="x", padx=10, pady=5)
        info_frame.grid_columnconfigure((0,1,2), weight=1)
        ctk.CTkLabel(info_frame, text=f"Top: {to}", font=("Arial", 14), text_color="#3498db").grid(row=0, column=0)
        ctk.CTkLabel(info_frame, text=f"Onay: {on}", font=("Arial", 14), text_color="#2ecc71").grid(row=0, column=1)
        ctk.CTkLabel(info_frame, text=f"Red: {re_val}", font=("Arial", 14), text_color="#e74c3c").grid(row=0, column=2)

    def rapora_yaz(self, metin, stil=None):
        if stil: self.txt_rapor._textbox.insert("end", metin, stil)
        else: self.txt_rapor.insert("end", metin)

    def formatla(self, isim, sayi1, sayi2=None, sayi3=None, sembol="▪"):
        isim_kisa = str(isim).upper().strip()[:45] 
        if sayi2 is not None and sayi3 is not None:
            return f" {sembol} {isim_kisa}\tToplam: {sayi1}\t| Onay: {sayi2}\t| Red: {sayi3}\n"
        else:
            return f" {sembol} {isim_kisa}\t: {sayi1}\n"

    def sutun_bul(self, df, aranan_kelimeler):
        for kelime in aranan_kelimeler:
            for col in df.columns:
                if kelime.lower() in str(col).lower(): return col
        return None

    # =========================================================
    # AKILLI RADAR: EN GÜNCEL DOSYAYI OTOMATİK BULUR
    # =========================================================
    def en_guncel_dosyayi_bul(self, anahtar_kelime):
        # İçinde anahtar kelime geçen ve sonu .xlsx ile biten tüm dosyaları bul
        dosyalar = glob.glob(f"*{anahtar_kelime}*.xlsx")
        # Gizli dosyaları veya geçici (~$) dosyaları ele
        dosyalar = [d for d in dosyalar if not d.startswith("~$")]
        if not dosyalar:
            return None
        # En son değiştirilme saatine göre en güncel olanı seç
        en_guncel_dosya = max(dosyalar, key=os.path.getmtime)
        return en_guncel_dosya

    def hesapla_onay_red_dinamik(self, df):
        if df.empty: return 0, 0, 0
        emniyet_col = self.sutun_bul(df, ["cagirma", "çağrı", "emniyet"])
        if not emniyet_col: return df.shape[0], 0, 0
        s_emniyet = df[emniyet_col].astype(str)
        onay = s_emniyet[s_emniyet.str.contains("onaylandı|onaylandi", case=False, na=False)].shape[0]
        red = s_emniyet[s_emniyet.str.contains("onaylanmadı|onaylanmadi", case=False, na=False)].shape[0]
        return df.shape[0], onay, red

    def dinamik_yukle(self):
        dosya_adi = filedialog.askopenfilename(
            title="Dinamik Rapor Excel dosyasını seçin",
            filetypes=[
                ("Excel Dosyaları", "*.xlsx *.xls"),
                ("Tüm Dosyalar", "*.*")
            ]
        )

        if not dosya_adi:
            messagebox.showwarning("İşlem İptal", "Dinamik Rapor dosyası seçilmedi.")
            return

        try:
            self.df_dinamik = pd.read_excel(dosya_adi)
            loc_col = self.sutun_bul(self.df_dinamik, ["lokasyon"])
            if loc_col:
                self.df_dinamik = self.df_dinamik[
                    self.df_dinamik[loc_col].astype(str).str.contains("dilovas", case=False, na=False)
                ]

            messagebox.showinfo(
                "Başarılı",
                f"Dinamik Rapor başarıyla yüklendi.\\n\\nOkunan Dosya:\\n{dosya_adi}"
            )

            if not self.df_pregate.empty:
                self.ana_sayfayi_goster()

        except Exception as e:
            messagebox.showerror("Hata", f"Dinamik Rapor dosyası okunamadı.\\n\\nDosya:\\n{dosya_adi}\\n\\nHata:\\n{e}")

    def pregate_yukle(self):
        dosya_adi = filedialog.askopenfilename(
            title="Pregate Sunum Excel dosyasını seçin",
            filetypes=[
                ("Excel Dosyaları", "*.xlsx *.xls"),
                ("Tüm Dosyalar", "*.*")
            ]
        )

        if not dosya_adi:
            messagebox.showwarning("İşlem İptal", "Pregate Sunum dosyası seçilmedi.")
            return

        try:
            self.df_pregate = pd.read_excel(dosya_adi)

            messagebox.showinfo(
                "Başarılı",
                f"Pregate Sunum başarıyla yüklendi.\\n\\nOkunan Dosya:\\n{dosya_adi}"
            )

            if not self.df_dinamik.empty:
                self.ana_sayfayi_goster()

        except Exception as e:
            messagebox.showerror("Hata", f"Pregate Sunum dosyası okunamadı.\\n\\nDosya:\\n{dosya_adi}\\n\\nHata:\\n{e}")

    def kpi_guncelle(self, toplam, onay, red):
        self.lbl_toplam.configure(text=str(toplam))
        self.lbl_onay.configure(text=str(onay))
        self.lbl_red.configure(text=str(red))

    def ekran_hazirla(self, baslik, goster_menu=False, mod="metin"):
        self.lbl_secili_dep.configure(text=baslik)
        self.txt_rapor.delete("1.0", "end")
        self.dashboard_frame.pack_forget()
        self.txt_rapor.pack_forget()
        self.tree_frame.pack_forget()
        
        if mod == "dashboard": self.dashboard_frame.pack(fill="both", expand=True, pady=10)
        elif mod == "tablo": self.tree_frame.pack(fill="both", expand=True, pady=10)
        else: self.txt_rapor.pack(pady=20, fill="both", expand=True)
            
        if goster_menu: self.nakliye_menu.pack(fill="x", pady=(0,10))
        else: self.nakliye_menu.pack_forget()

    def ana_sayfayi_goster(self):
        if self.df_dinamik.empty or self.df_pregate.empty:
            messagebox.showwarning("Uyarı", "Lütfen önce iki dosyayı da yükleyin.")
            return

        self.ekran_hazirla("🏠 GENEL OPERASYON DASHBOARD'U", mod="dashboard")
        for widget in self.dashboard_frame.winfo_children(): widget.destroy()

        havuz_col = self.sutun_bul(self.df_dinamik, ["havuz"])
        if havuz_col:
            s_havuz = self.df_dinamik[havuz_col].astype(str).str.lower().str.replace('i̇','i').str.replace('ı','i')
            kimya_mask = s_havuz.str.contains("kimya|yapi|yapı", na=False) & ~s_havuz.str.contains("dow", na=False)
            dow_mask = s_havuz.str.contains("dow", na=False)
            diger_mask = ~(kimya_mask | dow_mask)

            kimya_df = self.df_dinamik[kimya_mask]
            dow_df = self.df_dinamik[dow_mask]
            diger_df = self.df_dinamik[diger_mask]
        else: kimya_df, dow_df, diger_df = pd.DataFrame(), pd.DataFrame(), pd.DataFrame()

        k_t, k_o, k_r = self.hesapla_onay_red_dinamik(kimya_df)
        d_t, d_o, d_r = self.hesapla_onay_red_dinamik(dow_df)
        di_t, di_o, di_r = self.hesapla_onay_red_dinamik(diger_df)

        dep_col = self.sutun_bul(self.df_pregate, ["departman"])
        if dep_col:
            s_dep = self.df_pregate[dep_col].astype(str).str.lower()
            term_df = self.df_pregate[s_dep.str.contains("terminal", na=False)]
            antrepo_df = self.df_pregate[s_dep.str.contains("kuru|antrepo", na=False)]
        else: term_df, antrepo_df = pd.DataFrame(), pd.DataFrame()

        def pre_hesapla(df_temp):
            if df_temp.empty: return 0, 0, 0
            durum_col = self.sutun_bul(df_temp, ["durum"])
            iptal_col = self.sutun_bul(df_temp, ["eden", "iptal"])
            if not durum_col or not iptal_col: return df_temp.shape[0], 0, 0
            s_durum = df_temp[durum_col].astype(str).str.lower().str.strip()
            s_iptal = df_temp[iptal_col].astype(str).str.lower().str.strip()
            onay = s_durum[s_durum == "onay"].shape[0]
            p_red = s_durum.isin(["isg red", "ısg red", "sevkiyat red"]) & s_iptal.isin(PREGATE_EKIBI)
            s_red = (s_durum == "sevkiyat red") & ~s_iptal.isin(PREGATE_EKIBI)
            t_red = (s_durum == "terminal red")
            return df_temp.shape[0], onay, (p_red | s_red | t_red).sum()

        t_t, t_o, t_r = pre_hesapla(term_df)
        a_t, a_o, a_r = pre_hesapla(antrepo_df)

        self.kpi_guncelle(k_t+d_t+di_t+t_t+a_t, k_o+d_o+di_o+t_o+a_o, k_r+d_r+di_r+t_r+a_r)

        ctk.CTkLabel(self.dashboard_frame, text="Dinamik Rapor (Fabrika) Birimleri", font=("Arial", 20, "bold"), text_color="#E74C3C").grid(row=0, column=0, columnspan=3, pady=(10, 0))
        self.mini_karti_olustur(self.dashboard_frame, "Kimya Departmanı", k_t, k_o, k_r, 1, 0)
        self.mini_karti_olustur(self.dashboard_frame, "Dow Departmanı", d_t, d_o, d_r, 1, 1)
        self.mini_karti_olustur(self.dashboard_frame, "Diğer (Fabrika)", di_t, di_o, di_r, 1, 2)

        ctk.CTkLabel(self.dashboard_frame, text="Pregate Sunum (Terminal) Birimleri", font=("Arial", 20, "bold"), text_color="#3498DB").grid(row=2, column=0, columnspan=3, pady=(20, 0))
        self.mini_karti_olustur(self.dashboard_frame, "Poliport Terminal", t_t, t_o, t_r, 3, 0)
        self.mini_karti_olustur(self.dashboard_frame, "Antrepo (KuruYük)", a_t, a_o, a_r, 3, 1)

    def dinamik_rapor_analiz(self, departman):
        if self.df_dinamik.empty: return
        self.ekran_hazirla(f"Dinamik Rapor - {departman} Departmanı", mod="metin")
        
        df_temp = self.df_dinamik.copy()
        havuz_col = self.sutun_bul(df_temp, ["havuz"])
        
        if havuz_col:
            s_havuz = df_temp[havuz_col].astype(str).str.lower().str.replace('i̇','i').str.replace('ı','i')
            kimya_mask = s_havuz.str.contains("kimya|yapi|yapı", na=False) & ~s_havuz.str.contains("dow", na=False)
            dow_mask = s_havuz.str.contains("dow", na=False)
            
            if departman == "Kimya": df_temp = df_temp[kimya_mask]
            elif departman == "Dow": df_temp = df_temp[dow_mask]
            elif departman == "Diğer": df_temp = df_temp[~kimya_mask & ~dow_mask]
        
        toplam, onay, red = self.hesapla_onay_red_dinamik(df_temp)
        self.kpi_guncelle(toplam, onay, red)
        
        self.rapora_yaz(f"  {departman.upper()} DEPARTMAN GENEL ÖZETİ  \n", "baslik")
        self.rapora_yaz(self.formatla("Genel Toplam Kayıt", toplam, sembol="▶"), "satir_tek")
        self.rapora_yaz(self.formatla("Genel Onaylanan", onay, sembol="▶"), "satir_cift")
        self.rapora_yaz(self.formatla("Genel Reddedilen", red, sembol="▶"), "satir_tek")
        self.rapora_yaz("\n")
        
        if havuz_col and not df_temp.empty:
            self.rapora_yaz("  ALT BİRİM (HAVUZ) DETAYLARI  \n", "baslik")
            raw_pools = df_temp[havuz_col].dropna().unique()
            for i, pool in enumerate(sorted(raw_pools)):
                df_pool = df_temp[df_temp[havuz_col] == pool]
                p_t, p_o, p_r = self.hesapla_onay_red_dinamik(df_pool)
                satir = self.formatla(pool, p_t, p_o, p_r, sembol="▪")
                tag = "satir_tek" if i % 2 == 0 else "satir_cift"
                self.rapora_yaz(satir, tag)
            self.rapora_yaz("\n")
        
        self.rapora_yaz("  EN ÇOK ARAÇ GÖNDEREN NAKLİYECİLER (İLK 5)  \n", "baslik")
        nak_col = self.sutun_bul(df_temp, ["nakliyeci adi", "nakliyeci adı", "nakliye"])
        if nak_col:
            s_nak = df_temp[nak_col].astype(str).str.strip()
            top_nak = s_nak[s_nak != "nan"].value_counts().head(5)
            for i, (ad, sayi) in enumerate(top_nak.items()):
                satir = self.formatla(ad, f"{sayi} Araç", sembol="✔")
                tag = "satir_tek" if i % 2 == 0 else "satir_cift"
                self.rapora_yaz(satir, tag)
            self.rapora_yaz("\n")
            
        self.rapora_yaz("  KARŞILAŞILAN BAŞLICA RED NEDENLERİ (İLK 5)  \n", "baslik")
        neden_col = self.sutun_bul(df_temp, ["iptal neden", "neden", "aciklama"])
        if neden_col:
            s_neden = df_temp[neden_col].astype(str).str.strip()
            top_neden = s_neden[s_neden != "nan"].value_counts().head(5)
            for i, (neden, sayi) in enumerate(top_neden.items()):
                satir = self.formatla(neden, f"{sayi} Kez", sembol="✖")
                tag = "satir_tek" if i % 2 == 0 else "satir_cift"
                self.rapora_yaz(satir, tag)

    def pregate_analiz(self, departman):
        if self.df_pregate.empty: return
        self.ekran_hazirla(f"Pregate Sunum - {departman}", mod="metin")
        
        df_temp = self.df_pregate.copy()
        dep_col = self.sutun_bul(df_temp, ["departman"])
        if dep_col:
            s_dep = df_temp[dep_col].astype(str).str.lower()
            if departman == "Poliport Terminal": df_temp = df_temp[s_dep.str.contains("terminal", na=False)]
            elif departman == "Antrepo": df_temp = df_temp[s_dep.str.contains("kuru|antrepo", na=False)]

        durum_col = self.sutun_bul(df_temp, ["durum"])
        iptal_col = self.sutun_bul(df_temp, ["eden", "iptal"])

        if durum_col and iptal_col:
            s_durum = df_temp[durum_col].astype(str).str.lower().str.strip()
            s_iptal = df_temp[iptal_col].astype(str).str.lower().str.strip()
            
            p_red_mask = s_durum.isin(["isg red", "ısg red", "sevkiyat red"]) & s_iptal.isin(PREGATE_EKIBI)
            s_red_mask = (s_durum == "sevkiyat red") & ~s_iptal.isin(PREGATE_EKIBI)
            t_red_mask = (s_durum == "terminal red")
            
            red_df_all = df_temp[p_red_mask | s_red_mask | t_red_mask]
            
            toplam_arac = df_temp.shape[0]
            onay_sayisi = s_durum[s_durum == "onay"].shape[0]
            red_toplam = red_df_all.shape[0]
            
            self.kpi_guncelle(toplam_arac, onay_sayisi, red_toplam)
            
            self.rapora_yaz(f"  {departman.upper()} BİRİM BAZLI RED ANALİZİ  \n", "baslik")
            self.rapora_yaz(self.formatla("Pregate Personel Redleri (Sizin Ekip)", p_red_mask.sum(), sembol="▶"), "satir_tek")
            self.rapora_yaz(self.formatla("Terminal Sevkiyat Redleri (Diğer)", s_red_mask.sum(), sembol="▶"), "satir_cift")
            self.rapora_yaz(self.formatla("Terminal Operasyon Redleri", t_red_mask.sum(), sembol="▶"), "satir_tek")
            self.rapora_yaz("\n")

            self.rapora_yaz("  EN ÇOK RED EDİLEN NAKLİYECİLER (İLK 5)  \n", "baslik")
            nak_col = self.sutun_bul(red_df_all, ["nakliyeci adi", "nakliyeci adı", "nakliye"])
            if nak_col:
                s_nak = red_df_all[nak_col].astype(str).str.strip()
                top_nak = s_nak[s_nak != "nan"].value_counts().head(5)
                for i, (ad, sayi) in enumerate(top_nak.items()):
                    satir = self.formatla(ad, f"{sayi} Red", sembol="✖")
                    tag = "satir_tek" if i % 2 == 0 else "satir_cift"
                    self.rapora_yaz(satir, tag)
                self.rapora_yaz("\n")
            
            self.rapora_yaz("  BAŞLICA İPTAL AÇIKLAMALARI (İLK 5)  \n", "baslik")
            acik_col = self.sutun_bul(red_df_all, ["aciklama", "neden"])
            if acik_col:
                s_acik = red_df_all[acik_col].astype(str).str.strip()
                top_neden = s_acik[s_acik != "nan"].value_counts().head(5)
                for i, (neden, sayi) in enumerate(top_neden.items()):
                    satir = self.formatla(neden, f"{sayi} Kez", sembol="▪")
                    tag = "satir_tek" if i % 2 == 0 else "satir_cift"
                    self.rapora_yaz(satir, tag)

    def poliport_nakliyeci_arayuz(self):
        if self.df_pregate.empty:
            messagebox.showwarning("Uyarı", "Önce Pregate Sunum dosyasını yükleyin.")
            return
        
        self.ekran_hazirla("🚛 Poliport Terminal - Nakliyeci Analizi", goster_menu=True, mod="metin")
        df_temp = self.df_pregate.copy()
        dep_col = self.sutun_bul(df_temp, ["departman"])
        if dep_col:
            s_dep = df_temp[dep_col].astype(str).str.lower()
            self.df_pol_temp = df_temp[s_dep.str.contains("terminal", na=False)].copy()
        else: self.df_pol_temp = df_temp.copy()
            
        durum_col = self.sutun_bul(self.df_pol_temp, ["durum"])
        if durum_col:
            s_dur = self.df_pol_temp[durum_col].astype(str).str.lower().str.strip()
            t = self.df_pol_temp.shape[0]
            o = s_dur[s_dur == "onay"].shape[0]
            r = s_dur.isin(["isg red", "ısg red", "sevkiyat red", "terminal red"]).sum()
            self.kpi_guncelle(t, o, r)
            
        self.rapora_yaz("\n▶ Lütfen detaylarını görmek istediğiniz listeyi yukarıdaki butonlardan seçiniz.\n", "altbaslik")
        self.rapora_yaz("▶ 'Red Matrix Tablosu' butonuna tıklayarak Excel tablosu görünümüne geçebilirsiniz.\n", "altbaslik")

    def n_listele_tum(self):
        self.ekran_hazirla("🚛 Poliport Terminal - Tüm Nakliyeci Kayıtları", goster_menu=True, mod="metin")
        self.rapora_yaz("  TÜM NAKLİYECİ KAYIT BİLANÇOSU  \n", "baslik")
        for col in self.df_pol_temp.columns:
            if "nakliye" in str(col).lower():
                s_nak = self.df_pol_temp[col].astype(str).str.strip().str.upper()
                liste = s_nak[s_nak != "NAN"].value_counts()
                for i, (ad, sayi) in enumerate(liste.items()):
                    satir = self.formatla(ad, f"{sayi} Adet", sembol="▪")
                    tag = "satir_tek" if i % 2 == 0 else "satir_cift"
                    self.rapora_yaz(satir, tag)
                break

    def n_listele_onay(self):
        self.ekran_hazirla("🚛 Poliport Terminal - Onaylanan Nakliyeciler", goster_menu=True, mod="metin")
        self.rapora_yaz("  ONAYLANAN NAKLİYECİ KAYITLARI  \n", "baslik")
        durum_col = self.sutun_bul(self.df_pol_temp, ["durum"])
        if not durum_col: return
        onay_df = self.df_pol_temp[self.df_pol_temp[durum_col].astype(str).str.lower().str.strip() == "onay"]
        for col in onay_df.columns:
            if "nakliye" in str(col).lower():
                s_nak = onay_df[col].astype(str).str.strip().str.upper()
                liste = s_nak[s_nak != "NAN"].value_counts()
                for i, (ad, sayi) in enumerate(liste.items()):
                    satir = self.formatla(ad, f"{sayi} Onay", sembol="✔")
                    tag = "satir_tek" if i % 2 == 0 else "satir_cift"
                    self.rapora_yaz(satir, tag)
                break

    def yapay_zeka_kategori(self, metin):
        m = str(metin).lower().replace('i̇','i').replace('ı','i').replace('ş','s').replace('ğ','g').replace('ü','u').replace('ö','o').replace('ç','c')
        if any(k in m for k in ["sorumluluk", "taahhutname"]): return "Sorumluluk Evrakı"
        if any(k in m for k in ["yuksek", "saglik", "rapor"]): return "Yüksekte Çalışabilir Sağlık Raporu"
        if any(k in m for k in ["basinc", "t9"]): return "Tank Basınç Raporu - T9"
        if any(k in m for k in ["muayene", "vize"]): return "Muayene"
        if any(k in m for k in ["trafik", "sigorta"]) and not "atik" in m: return "Zorunlu Trafik Sigortası"
        if any(k in m for k in ["k1"]): return "K1 Taşıt Kartı"
        if any(k in m for k in ["tehlikeli atik"]): return "Tehlikeli Atık Sigortası"
        if any(k in m for k in ["tank kodu"]): return "Tank Kodu"
        if any(k in m for k in ["kkd", "baret", "gozluk", "kemer", "reflektor", "ayakkabi"]): return "KKD"
        if any(k in m for k in ["interchange"]): return "interchange"
        if any(k in m for k in ["kart almaya gelmedi", "kart almadi"]): return "Şoför Kart Almaya Gelmedi"
        if any(k in m for k in ["firma iptal", "surucu iptal", "sevkiyat iptal", "gelmedi", "arac arizasi"]): return "Sevkiyat - Firma Sürücü İptal"
        if any(k in m for k in ["havuz", "tahliye vanasi", "vana"]): return "Havuz Tahliye Vanası"
        if any(k in m for k in ["isopa"]): return "İSOPA"
        if any(k in m for k in ["far", "stop", "sis", "lastik", "kirik", "patlak", "aydinlatma", "lamba", "fari"]): return "FAR-STOP-SİS FARI-LASTİK KIRIK/PATLAK"
        if any(k in m for k in ["turuncu plaka", "ikaz", "tabela", "levha", "plaka"]): return "Turuncu Plaka ve İkaz Tabelası"
        if any(k in m for k in ["sizinti", "kacak", "yakit deposu", "yakit"]): return "Yakıt Deposu Sızıntı"
        return "Diğer"

    def get_matrix_dataframe(self):
        durum_col = self.sutun_bul(self.df_pol_temp, ["durum"])
        nak_col = self.sutun_bul(self.df_pol_temp, ["nakliyeci adi", "nakliye"])
        acik_col = self.sutun_bul(self.df_pol_temp, ["aciklama", "neden"])
        if not durum_col or not nak_col or not acik_col: return pd.DataFrame()
        
        df_temp = self.df_pol_temp.copy()
        df_temp["Temiz_Nakliyeci"] = df_temp[nak_col].astype(str).str.strip().str.upper()
        toplam_kayit = df_temp[df_temp["Temiz_Nakliyeci"] != "NAN"]["Temiz_Nakliyeci"].value_counts()
        
        s_durum = df_temp[durum_col].astype(str).str.lower().str.strip()
        red_df = df_temp[s_durum.isin(["isg red", "ısg red", "sevkiyat red", "terminal red"])].copy()
        
        if red_df.empty: return pd.DataFrame()
        
        red_df["Kategori"] = red_df[acik_col].apply(self.yapay_zeka_kategori)
        pivot = pd.crosstab(red_df["Temiz_Nakliyeci"], red_df["Kategori"])
        
        for k in KATEGORILER:
            if k not in pivot.columns: pivot[k] = 0
                
        pivot = pivot[KATEGORILER]
        pivot["Toplam Red"] = pivot.sum(axis=1)
        pivot["Toplam Kayıt"] = pivot.index.map(toplam_kayit).fillna(0)
        pivot = pivot.sort_values(by="Toplam Red", ascending=False)
        pivot.loc["TOPLAM"] = pivot.sum()
        return pivot

    def n_listele_red(self):
        self.ekran_hazirla("🚛 Poliport Terminal - Nakliyeci Red Matrix Dağılımı", goster_menu=True, mod="tablo")
        self.tree.delete(*self.tree.get_children())
        
        pivot = self.get_matrix_dataframe()
        if pivot.empty: return
        
        cols = ["Firma Adı"] + list(pivot.columns)
        self.tree["columns"] = cols
        self.tree.column("#0", width=0, stretch="no")
        self.tree.heading("#0", text="")

        for col in cols:
            w = 250 if col == "Firma Adı" else 150
            if col in ["Toplam Red", "Toplam Kayıt"]: w = 120
            self.tree.column(col, anchor="center", width=w)
            self.tree.heading(col, text=col, anchor="center")

        for index, row in pivot.iterrows():
            vals = [index] + [int(row[c]) if row[c] != 0 else "" for c in pivot.columns]
            tag = "toplam" if index == "TOPLAM" else "normal"
            self.tree.insert("", "end", values=vals, tags=(tag,))

        self.tree.tag_configure("toplam", background="#1F242D", foreground="#F1C40F", font=("Arial", 12, "bold"))
        self.tree.tag_configure("normal", font=("Arial", 11))

    # =========================================================
    # EXPORT (DIŞA AKTARMA) ALANI
    # =========================================================
    def export_excel(self):
        if self.df_dinamik.empty or self.df_pregate.empty:
            messagebox.showwarning("Uyarı", "Lütfen önce dosyaları yükleyin.")
            return
        
        try:
            writer = pd.ExcelWriter("Operasyon_Detayli_Rapor.xlsx", engine="openpyxl")
            pivot = self.get_matrix_dataframe()
            if not pivot.empty:
                pivot.reset_index(names="Firma Adı").to_excel(writer, sheet_name="Nakliyeci Red Matrix", index=False)
            else:
                pd.DataFrame({"Bilgi": ["Matrix için veri bulunamadı."]}).to_excel(writer, sheet_name="Nakliyeci Red Matrix")
                
            writer.close()
            messagebox.showinfo("Başarılı", "Operasyon_Detayli_Rapor.xlsx dosyası oluşturuldu.")
        except Exception as e:
            messagebox.showerror("Hata", f"Excel oluşturulamadı: {e}")

    def export_pptx(self):
        if self.df_dinamik.empty or self.df_pregate.empty:
            messagebox.showwarning("Uyarı", "Lütfen önce dosyaları yükleyin.")
            return
            
        try:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Pregate Operasyon Raporu"
            slide.placeholders[1].text = "Hazırlayan: S. SEYMEN"
            
            slide2 = prs.slides.add_slide(prs.slide_layouts[1])
            slide2.shapes.title.text = "Fabrika (Dinamik) Özet"
            tf = slide2.placeholders[1].text_frame
            tf.text = "Kimya ve Dow birimlerindeki güncel iş akışı detayları programdan incelenebilir."
            
            prs.save("Operasyon_Sunumu.pptx")
            messagebox.showinfo("Başarılı", "Operasyon_Sunumu.pptx başarıyla oluşturuldu.")
        except Exception as e:
            messagebox.showerror("Hata", f"PowerPoint oluşturulamadı: {e}\n(Terminal'e 'pip3 install python-pptx' yazdığınızdan emin olun.)")

if __name__ == "__main__":
    app = PregateDashboard()
    app.mainloop()